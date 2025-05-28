# 환경 변수 적재하기
from dotenv import load_dotenv

load_dotenv()

# OpenAI 인스턴스 생성
from openai import OpenAI

client = OpenAI()

import pandas as pd

# 파일 경로 설정
file_path = 'data/sales_data.csv'
sales_data = pd.read_csv(file_path)

sales_data.head()

# 파일 읽기
file = client.files.create(
    file=open(file_path, 'rb'),
    purpose='assistants',
)

# 화일을 포함한 assistant 생성
assistant = client.beta.assistants.create(
    instructions='데이터 과학 도우미로서, 주어진 데이터와 요청에 따라 적절한 코드를 작성하고 적절한 시각화를 생성 할 수 있습니다.',
    model='gpt-4o-mini',
    tools=[
        { 'type': 'code_interpreter' }
    ],
    tool_resources={
        'code_interpreter': {
            'file_ids': [file.id]
        }
    }
)

print('--- assistant 생성 완료 ---')
print(assistant)

# 스레드 생성
thread = client.beta.threads.create(
    messages=[
        {
            'role': 'user',
            'content': '2022년부터 2025년까지 각 분기의 총 판매액을 계산하고, 이를 다른 제품으로 시각화하여 선 그래프로 표시하세요. 제품의 선 색상은 각각 빨강, 파랑, 녹색으로 설정하세요.',
            'attachments': [
                {
                    'file_id': file.id,
                    'tools': [
                        { 'type': 'code_interpreter' }
                    ]
                }
            ]
        }
    ]
)

print('--- 스레드 생성 완료 ---')
print(thread)

# 실행 세션 생성
run = client.beta.threads.runs.create(
    thread_id=thread.id,
    assistant_id=assistant.id,
)

# 실행 세션 출력
print('--- 실행 세션 생성 완료 ---')
print(run)

import time

while True:
    messages = client.beta.threads.messages.list(
        thread_id=thread.id,
    )
    
    try:
        if messages.data and messages.data[0].content and messages.data[0].content[0].image_file:
            print('--- 이미지 chart 생성 완료 ---')
            if messages.data and messages.data[0].content:
                print('현재 메시지:', messages.data[0].content[0])
            break
    except AttributeError:
        time.sleep(10)
        print('--- 도우미가 차트를 작성 하는 중.... (AttributeError)')
        if messages.data and messages.data[0].content:
            print('현재 메시지:', messages.data[0].content[0])
    except IndexError: # Handles cases where data or content list is empty
            time.sleep(10)
            print('--- 도우미가 차트를 작성 하는 중.... (IndexError)')
            # Optional: print a message if messages.data is empty or content is empty
            if not messages.data:
                print("메시지 목록이 비어 있습니다.")
            elif not messages.data[0].content:
                print("첫 번째 메시지에 내용이 없습니다.")
    # 순환 전 잠시 대개
    time.sleep(5)
    
# 화일을 PNG 형식으로 변환하는 함수
def convert_file_to_png(file_id, write_path):
    data = client.files.content(file_id)
    data_bytes = data.read()
    
    with open(write_path, 'wb') as file:
        file.write(data_bytes)
        
#---------------------------

# 첫번째 메시지에서 이미지 화일 ID 가져오기
plot_file_id = messages.data[0].content[0].image_file.file_id
image_path = '태진_도서판매.png'

# 화일을 PNG로 변환
convert_file_to_png(plot_file_id, image_path)

# chart 업로드
plot_file = client.files.create(
    file=open(image_path, 'rb'),
    purpose='assistants',
)

# 도우미의 생각과 행동 과정 표시
messages = client.beta.threads.messages.list(thread_id=thread.id)
assistant_thoughts_and_actions = [message.content[0] for message in messages.data]

# 결과 출력
for content in assistant_thoughts_and_actions:
    print(content)
    
import time

# 사용자 메시지를 제출하고 완료를 기다리는 함수
def submit_message_wait_completion(assistant_id, thread, user_message, file_ids=None):
    # 활성화된 실행 세션이 완료 될 떄까지 대기
    for run in client.beta.threads.runs.list(thread_id=thread.id).data:
        if run.status == 'in_progress':
            print(f'실행 세션 {run.id} 완료 대기 중...')
            
            while True:
                run_status = client.beta.threads.runs.retrieve(thread_id=thread.id, run_id=run.id).status
                
                if run_status in ['succeeded', 'failed']:
                    break
                
                time.sleep(5)
                
    # 메시지 제출
    params = {
        'thread_id': thread.id,
        'role': 'user',
        'content': user_message,
    }
    
    # 첨부화일 설정
    if file_ids:
        attachments = [{"file_id": file_id, "tools": [{"type": "code_interpreter"}]} for file_id in file_ids]
        params['attachments'] = attachments
        
    _ = client.beta.threads.messages.create(**params)
    
    # 실행 세션 생성
    run = client.beta.threads.runs.create(
        thread_id=thread.id,
        assistant_id=assistant_id,
    )
    return run

#---------------------------------------------------
# 요청을 보내 도우미에게 통찰 생성을 요청
submit_message_wait_completion(
    assistant_id=assistant.id,
    thread=thread,
    user_message='앞에서 생성한 차트를 기반으로 약 20자 내외의 문장 두 개로 가장 중요한 통찰을 설명해 주세요. 이 내용은 프레젠테이션 발표에서 데이터의 비밀을 드러내기 위해 사용될 것입니다.'
)
    
# 대화 흐름의 응답을 가져오는 함수
def get_response(thread):
    return client.beta.threads.messages.list(thread_id=thread.id)

#---------------------------------------------------

# 응답 대기후 생성된 통찰 가져오기
time.sleep(10)

response = get_response(thread)
bullet_points = response.data[0].content[0].text.value

print('--- 통찰 가져오기 완료 ---')
print(bullet_points)

#---------------------------------------------------

# 통찰에 기반한 제목 생성
submit_message_wait_completion(
    assistant_id=assistant.id,
    thread=thread,
    user_message='당신이 만든 차트와 통찰을 바탕으로, 주요 통찰을 반영하는 아주 짧은 프레젠테이션 제목을 만들어 주세요.'
)

# 응답 대기 후 생성된 제목 출력
time.sleep(10)

response = get_response(thread)
title = response.data[0].content[0].text.value

print('--- 제목 생성 완료 ---')
print(title)

#---------------------------------------------------
# 회사 설명 제공
company_summary = '비록 우리는 신생 온라인 꽃 도매 전자상거래 회사에 불과하지만, 회장님은 IT 서적도 집필합니다!'

# DALL-E 3 모델을 호출하여 이미지 생성
response = client.images.generate(
    model='dall-e-3',
    prompt=f'이 회사의 설명인 {company_summary}을 바탕으로, 태진과 꽃말의 비밀정원 회사가 함께 성장하고 전진하는 영감을 주는 이미지를 만들어 주세요. 이 이미지는 분기별 판매 계획 회의에서 사용될 것입니다.',
    size='1024x1024',
    quality='hd',
    n=1,
)

# DALL-E 3 모델이 생성한 이미지 URL 가져오기
image_url = response.data[0].url

# DALL-E 3 모델이 생성한 이미지 가져오기 및 저장
import requests

dalle_image_path = '꽃말의비밀정원_태진.png'
img = requests.get(image_url)

# 생성된 이미지 저장
with open(dalle_image_path, 'wb') as file:
    file.write(img.content)

# 생성된 이미지를 프레젠테이션 자료로 사용
dalle_file = client.files.create(
    file=open(dalle_image_path, 'rb'),
    purpose='assistants',
)

#---------------------------------------------------
# 프레젠테이션 프레임 생성
# 첫 페이지 템플릿
title_template = """
# 모듈 적재
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# 프레젠테이션 객체 생성
prs = Presentation()

# 레이아웃 추가
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 배경색 설정
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# 이미지 추가
left = Inches(0)
top = Inches(0)
height = prs.slide_height
width = prs.slide_width * 3 / 5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# 제목 텍스트 상자
left = prs.slide_width * 3 / 5
top = Inches(2)
width = prs.slide_width * 2 / 5
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(38)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# 부제목 텍스트 상자
left = prs.slide_width * 3 / 5
top = Inches(3)
width = prs.slide_width * 2 / 5
height = Inches(1)
subtitle_box = slide.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = subtitle_text
subtitle_p.font.size = Pt(22)
subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
subtitle_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
"""

# 나머지 페이지 템플릿
details_template = """
# 프레젠테이션 객체 생성
prs = Presentation()

# 레이아웃 추가
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 배경색 설정
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# 이미지 경로와 텍스트 정의
image_path = data_vis_img
title_text = "향상된 수익성: 온라인 판매와 직접 판매 최적화의 주도적 역할"
bullet_points = (
    "• 온라인 판매는 분기마다 수익성에서 선두를 유지하며, 강력한 디지털 시장이 존재함을 나타냅니다.\n"
    "• 직접 판매는 변동성이 있으며, 이 채널에서 성과 변화와 목표 지향적 개선의 필요성을 보여줍니다."
)

# 이미지 추가
left = Inches(0.2)
top = Inches(1.8)
height = prs.slide_height - Inches(3)
width = prs.slide_width * 3 / 5
pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

# 제목 추가
left = Inches(0)
top = Inches(0)
width = prs.slide_width
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.margin_top = Inches(0.1)
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(28)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# 주요 통찰 텍스트와 글머리 기호 목록 추가
left = prs.slide_width * 2 / 3
top = Inches(1.5)
width = prs.slide_width * 1 / 3
height = Inches(4.5)
insights_box = slide.shapes.add_textbox(left, top, width, height)
insights_frame = insights_box.text_frame

# 주요 통찰 텍스트 추가
insights_p = insights_frame.add_paragraph()
insights_p.text = "주요 통찰:"
insights_p.font.bold = True
insights_p.font.size = Pt(24)
insights_p.font.color.rgb = RGBColor(0, 128, 100)
insights_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# 글머리 기호 목록 추가
bullet_p = insights_frame.add_paragraph()
bullet_p.text = bullet_points
bullet_p.font.size = Pt(12)
bullet_p.font.color.rgb = RGBColor(255, 255, 255)
bullet_p.line_spacing = 1.5
"""

#---------------------------------------------------

dalle_img_path = '꽃말의비밀정원_태진.png'
dalle_file = client.files.create(
    file=open(dalle_img_path, 'rb'),
    purpose='assistants',
)

image_path = '태진_도서판매.png'
plot_file = client.files.create(
    file=open(image_path, 'rb'),
    purpose='assistants',
)

title_text = '꽃말의 비밀정원'
subtitle_text = '2025년 판매 회의'

# 프레젠테이션 생성 요청
submit_message_wait_completion(
    assistant_id=assistant.id,
    thread=thread,
    user_message=(
        f"포함된 코드 템플릿을 사용하여 템플릿 형식에 맞는 프레젠테이션을 생성하세요. 이 메시지에 포함된 이미지, 회사 이름/제목 및 파일명/부제목을 사용하세요:\n"
        f"{title_template}.\n"
        f"첫 번째 페이지에는 이 메시지에 포함된 이미지 파일 dalle_file.id를 image_path로 사용하고, 회사 이름 '{title_text}'을 title_text 변수로 사용하고, "
        f"부제목 텍스트 '{subtitle_text}'을 subtitle_text 변수로 사용하세요.\n"
        f"두 번째 페이지에는 다음 코드 템플릿 {details_template}을 사용하여 템플릿 형식에 맞는 프레젠테이션을 생성하세요.\n"
        f"중요: 두 번째 첨부 이미지인 꺾은선 그래프 plot_file.id를 data_vis_img로 사용하고, 이전에 생성한 데이터 시각화 제목을 title_text로 사용하며, "
        f"이전에 생성한 인사이트 글머리 기호 목록을 bullet_points 변수로 사용하세요."
        f"이 두 페이지를 pptx 형식의 파일로 출력하세요. 각각의 페이지가 이 메시지에서 제공된 템플릿 형식에 맞도록 해야 합니다."
    ),
    file_ids=[dalle_file.id, plot_file.id],
)

# 생성 작업 완료 대기
while True:
    try:
        response = get_response(thread)
        pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id
        print(f'--- 프레젠테이션 생성 완료: 성공적으로 pptx_id를 찾았습니다. ---')
        break
    except Exception as e:
        print('도우미가 슬라이드를 열심히 제작하고 있습니다...')
        print(response.data[0].content[0].text)
        time.sleep(10)
        
# 생성된 프레젠테이션 pptx 파일 다운로드/저장
import io

pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id

pptx_file = client.files.content(pptx_id)
pptx_bytes = io.BytesIO(pptx_file.read())

with open('태진_꽃말의비밀정원.pptx', 'wb') as f:
    f.write(pptx_bytes.getbuffer())
    
print('--- 프레젠테이션 pptx 파일 저장 완료 ---')
print('----------------------------------------------')


